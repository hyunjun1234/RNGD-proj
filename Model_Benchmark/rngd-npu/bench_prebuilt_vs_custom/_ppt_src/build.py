#!/usr/bin/env python3
"""prebuilt 대 custom 모델 벤치 덱. 수치는 전부 ../summary.json 에서 읽는다(손으로 적지 않는다).

서사(L-39): 무엇을 비교하나 → 어떻게 쟀나(코드) → 속도 → 적재 → 답변 → 원인(코드) → 검증.
머리글은 "그림 N: 무엇을 그렸나"(L-38), 코드 캡션은 "코드 N: 파일, 무엇".
"""
import sys, os, json
sys.path.insert(0, "/home/jun/.claude/skills/diagram-deck/scripts")
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deck
from check_shapes import compare, check_group_tight
SRC = os.path.dirname(os.path.abspath(__file__))
SUM = json.load(open(os.path.join(SRC, "..", "summary.json"), encoding="utf-8"))
BY = {d["model"]: d for d in SUM}


def g(mid, *path, default=None):
    d = BY.get(mid)
    for p in path:
        if not isinstance(d, dict):
            return default
        d = d.get(p)
    return default if d is None else d


def load(n):
    return open(os.path.join(SRC, n + ".svg"), encoding="utf-8").read()


d = deck.Deck("직접 빌드한 모델은 배포판보다 빠른가",
              "같은 입력, 같은 조건으로 prebuilt FXB 와 우리가 빌드한 tp8 아티팩트를 재 봤다",
              "RNGD 4장 서버, 2026-08-28 실측")
srcs = {}


def dia(title, name, head, items):
    svg = load(name)
    s, n, mp, _ = d.diagram_svg(title, svg, items, None, None, None,
                                name="도해 " + name, cols=1, head=head)
    srcs["도해 " + name] = svg
    print(f"{name}: 도형 {n}개, 최소 글자 {mp:.1f}pt")


def build():
    dia("무엇과 무엇을 비교했나", "paths",
        head="그림 1: 같은 모델이 서빙되는 두 경로, 배포 FXB 번들과 우리가 빌드한 v2 아티팩트",
        items=[
            {"t": "라우터가 같은 모델을 두 갈래로 띄운다. 이름 그대로면 퓨리오사가 배포한 FXB 번들(tp32, 카드 4장), 이름 뒤에 @tp8 을 붙이면 우리가 2026-07-27~29 에 직접 빌드한 v2 아티팩트(tp8, 카드 1장)다"},
            {"t": "두 갈래는 런타임 코드 경로부터 다르다. serve 로그를 보면 FXB 는 furiosa_generator::entrypoint 를, v2 아티팩트는 furiosa_generator::next_gen 을 탄다. 같은 가중치라도 실행 경로가 같지 않다"},
            {"t": "직접 빌드한 8개 중 5개는 MoE 게이트를 지나려고 artifact.json 의 model_type 을 qwen3_moe 에서 qwen3 로 고쳐 두었다. 나머지 3개(Qwen3-32B, EXAONE, Llama)는 손대지 않았다"},
        ])
    d.code("측정 방법", [
        "# bench.py — 모델 하나마다 (1) 준비될 때까지 기다리고 (2) 프롬프트 4종을 스트리밍으로 재고",
        "#            (3) 같은 프롬프트 4개를 동시에 던져 총처리량을 잰다",
        "PROMPTS   = [사실, 코드, 추론, 설명]        # 네 종류 고정, 모든 모델에 같은 문장",
        "MAX_TOKENS = 1024                          # 256 이면 thinking 모델이 사고만 하다 끝난다",
        "body = {'model': model, 'messages': [...], 'temperature': 0,   # greedy, 무작위성 제거",
        "        'max_tokens': MAX_TOKENS, 'stream': True,",
        "        'stream_options': {'include_usage': True}}             # 토큰 수를 서버에서 받는다",
        "",
        "t0 = time.perf_counter()",
        "for 이벤트 in 스트림:",
        "    조각 = delta.content 또는 delta.reasoning        # thinking 모델은 reasoning 으로 온다",
        "    if 조각 and ttft is None:",
        "        ttft = time.perf_counter() - t0              # 첫 토큰까지 = TTFT",
        "total = time.perf_counter() - t0",
        "decode_tps = (출력토큰 - 1) / (total - ttft)         # 첫 토큰 뒤의 순수 생성 속도",
        "",
        "# 준비 판정: 상태가 'up' 이 되기 전에 요청을 보내면 타임아웃까지 통째로 버린다(실측 600s)",
        "ready = router_status()['running'].get(model, {}).get('state') == 'up'",
    ], caption="코드 1: bench.py, 같은 조건을 강제하는 부분과 시간을 재는 부분")
    dia("무엇을 쟀나", "method",
        head="그림 2: 한 요청의 시간 축에서 TTFT 와 디코드 구간이 각각 어디인지, 그리고 같게 맞춘 조건 넷",
        items=[
            {"t": "TTFT 는 요청을 보낸 순간부터 첫 토큰이 올 때까지다. 프롬프트를 읽는 시간(prefill)이 여기 들어간다. 디코드 속도는 첫 토큰 이후 남은 토큰을 초로 나눈 값이라 프롬프트 길이에 영향을 덜 받는다"},
            {"t": "총처리량은 같은 프롬프트 4개를 동시에 던져 전부 끝날 때까지의 벽시계로 나눈 값이다. 배치가 되는 서버일수록 단일 요청보다 크게 나온다"},
            {"t": "적재 시간은 bench.py 가 잰 값(전환 시간, 앞 모델을 내리는 시간 포함)과 serve 로그에서 뽑은 순수 적재 시간 두 가지를 따로 본다"},
        ])
    dia("속도는 어느 쪽이 빠른가", "speed",
        head="그림 3: 모델 쌍마다 prebuilt 와 custom 의 동시 4요청 총처리량, 그리고 단일 요청 디코드 속도",
        items=[
            {"t": "카드 수가 다르다는 점을 먼저 봐야 한다. prebuilt 는 4장을 독점하고 custom tp8 은 1장만 쓴다. 카드당으로 환산하면 격차가 네 배 더 벌어진다"},
            {"t": "TTFT 는 두 갈래 모두 0.1초 안팎이라 사람이 느끼는 첫 응답 속도에는 차이가 없다. 차이는 디코드 속도와 동시 처리에서 난다"},
            {"t": "다만 이 숫자는 답변이 멀쩡한 경우에만 뜻이 있다. 깨진 모델도 토큰은 빠르게 뱉는다"},
        ])
    dia("모델을 올리는 데 얼마나 걸리나", "load",
        head="그림 4: serve 로그 기준 순수 적재 시간, 배포 FXB 와 직접 빌드한 tp8 아티팩트",
        items=[
            {"t": "라우터는 카드가 모자라면 쓰던 모델을 내리고 새로 올린다. 그래서 사용자가 체감하는 전환 시간은 여기에 내리는 시간이 더해진다"},
            {"t": "tp8 아티팩트는 카드 하나에 올리므로 짧고, tp32 FXB 는 4장에 펼치느라 길다. 모델을 자주 바꾸는 환경이라면 이 차이가 크게 다가온다"},
        ])
    dia("답변은 멀쩡한가", "quality",
        head="그림 5: 직접 빌드한 아티팩트별로 네 프롬프트의 답변이 정상인지, MoE 위장 여부와 나란히",
        items=[
            {"t": "속도만 보면 custom 이 이기는 것처럼 보이지만, 답변을 같이 보면 MoE 위장을 한 아티팩트가 뜻 없는 문자를 반복한다. 처리량 지표는 깨진 출력도 똑같이 빠르게 센다"},
            {"t": "위장하지 않은 아티팩트(Qwen3-32B, EXAONE, Llama)는 셋 다 정상이고, 위장한 다섯은 모두 깨졌다. 깨진 쪽은 같은 문자를 반복하거나 질문과 무관한 다른 언어를 쏟아낸다"},
        ])
    d.code("게이트를 우회한 코드와 그 대가", [
        "# masquerade_moe.sh 가 한 일: artifact.json 의 model_type 문자열을 바꾼다",
        '  "model_metadata": {',
        '      "model_type": "qwen3_moe",        →  "qwen3"      # 게이트가 보는 값',
        '      "hf_configs": {',
        '          "architectures": ["Qwen3MoeForCausalLM"]  →  ["Qwen3ForCausalLM"]',
        '          "num_local_experts": 128,     ← 이 MoE 키 8개가 함께 사라졌다',
        '          "num_experts_per_tok": 8,',
        '          "moe_intermediate_size": 768, …',
        "",
        "# 게이트가 실제로 뱉는 말 (원본 그대로 띄웠을 때)",
        "pyo3_runtime.PanicException: Unsupported model metadata:",
        "    ModelMetadata { model_type: Some(Qwen3Moe), task: Some(Generate), … }",
        "",
        "# 검증 실험(moe_check2.py): 세 가지로 바꿔 가며 같은 질문을 던진다",
        "for 변형 in [위장본, 원본 qwen3_moe, 위장 + MoE키 8개 복원]:",
        "    artifact.json 을 그 변형으로 교체",
        "    kill_backend(model)          # 라우터가 다시 올리도록 백엔드를 죽인다",
        "    답 = 라우터에 같은 질문",
        "# 결과: 위장본 깨짐 / 원본 안 뜸 / MoE키 복원본도 깨짐 → 키 소실은 원인이 아니다",
    ], caption="코드 2: masquerade_moe.sh 가 바꾼 필드와 moe_check2.py 의 검증 절차")
    dia("게이트는 무엇을 막고 있었나", "gate",
        head="그림 6: 같은 MoE 모델이 배포 FXB 와 직접 빌드 v2 에서 갈리는 지점, 그리고 우회했을 때 벌어지는 일",
        items=[
            {"t": "furiosa-llm 2026.3.0 의 v2 아티팩트 경로는 Qwen3Moe 를 지원하지 않는다고 명시적으로 거부한다. 같은 MoE 모델이라도 퓨리오사가 배포한 FXB 번들로는 정상 동작한다 — 이번 측정에서 prebuilt Coder 와 A3B 계열이 모두 제대로 답했다"},
            {"t": "model_type 문자열만 바꾸면 검사는 지나가지만 실행은 고쳐지지 않는다. 가중치는 MoE 구조 그대로인데 런타임이 dense 로 다루므로 결과가 틀린다"},
            {"t": "이 실패는 예외도 경고도 남기지 않는다. 오히려 처리량과 지연은 정상 모델과 비슷하거나 더 좋게 나오므로, 속도만 재는 벤치는 이것을 걸러내지 못한다"},
        ])
    dia("사라진 설정 키가 원인인가", "verify",
        head="그림 7: artifact.json 을 세 가지로 바꿔 같은 질문을 던진 결과, 가설을 기각한 실험",
        items=[
            {"t": "위장할 때 MoE 설정 키 8개가 함께 지워지는 것은 사실이다. 그래서 그 키만 되살리고 model_type 위장은 유지한 절충본을 만들어 다시 띄워 봤다"},
            {"t": "절충본도 여전히 질문에 답하지 못했다. 따라서 원인은 메타데이터에 빠진 값이 아니라 실행 경로가 MoE 를 다루지 못한다는 사실 자체다"},
            {"t": "이 실험은 라우터를 거쳐 실제 서빙 경로로 했다. 끝난 뒤 artifact.json 은 원래 상태로 되돌려 두었다"},
        ])


if __name__ == "__main__":
    build()
    OUT = os.path.join(SRC, "..", "prebuilt-vs-custom.pptx")
    d.save(OUT)
    print("넘침:", deck.OVERFLOW if deck.OVERFLOW else "없음")
    from pptx import Presentation
    prs = Presentation(OUT); worst = 0.0
    for sl in prs.slides:
        for shp in sl.shapes:
            if shp.shape_type == 6 and shp.name in srcs:
                diff, _, _ = compare(srcs[shp.name], shp); worst = max(worst, diff)
                print(f"되돌려 비교 {shp.name}: {diff*100:.2f}% ({'통과' if diff <= 0.02 else '실패'})")
    print(f"최대 {worst*100:.2f}%, 슬라이드 {len(prs.slides)}장, 저장 {os.path.abspath(OUT)}")
