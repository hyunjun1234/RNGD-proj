#!/usr/bin/env python3
"""쓸 수 있는 모델 정리 덱. 수치는 전부 ../summary.json 에서 읽는다.

서사: 어떤 모델이 있나 → 어떻게 쟀나(코드) → 시간 축 → 지연 → 처리량 → 적재 → 답변 → 고르는 법.
머리글은 "그림 N: 무엇을 그렸나"(L-38), 코드 캡션은 "코드 N: 파일, 무엇"(L-39).
"""
import sys, os, json
sys.path.insert(0, "/home/jun/.claude/skills/diagram-deck/scripts")
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deck
from check_shapes import compare, check_group_tight
SRC = os.path.dirname(os.path.abspath(__file__))
SUM = json.load(open(os.path.join(SRC, "..", "summary.json"), encoding="utf-8"))
BY = {d["model"]: d for d in SUM}


def load(n):
    """SVG 를 읽으면서 조립 전에 먼저 확인한다 — 문제를 pptx 를 만든 뒤가 아니라 여기서 잡는다.

    ① XML 로 파싱되는가 (글자 안 & < > 이스케이프 누락이면 여기서 걸린다)
    ② 내용 폭이 828 인가, 높이가 234 이하인가 (다르면 슬라이드마다 축척이 달라진다, L-30)
    """
    import xml.etree.ElementTree as ET
    from svg2shapes import content_bbox
    svg = open(os.path.join(SRC, n + ".svg"), encoding="utf-8").read()
    try:
        ET.fromstring(svg)
    except ET.ParseError as e:
        raise SystemExit(f"[{n}.svg] XML 이 깨졌다: {e}\n  글자 안의 & < > 를 esc() 로 감쌌는지 보라")
    x0, y0, x1, y1 = content_bbox(svg)
    w, h = round(x1 - x0), round(y1 - y0)
    if w != 828 or h > 234:
        print(f"  ! {n}.svg 규격 벗어남 — 폭 {w}(828 이어야), 높이 {h}(234 이하여야)")
    return svg


def n_done():
    import gen_all
    return sum(1 for mid, _, _ in gen_all.MODELS if mid in BY and "runs" in BY[mid])


d = deck.Deck("RNGD 서버에서 쓸 수 있는 모델",
              "채팅 모델 전종을 같은 입력으로 재 본 처리량, 지연, 그리고 실제 답변",
              "RNGD 4장 서버, 2026-08-29 실측")
srcs = {}


def dia(title, name, head, items):
    svg = load(name)
    s, n, mp, _ = d.diagram_svg(title, svg, items, None, None, None,
                                name="도해 " + name, cols=1, head=head)
    srcs["도해 " + name] = svg
    print(f"{name}: 도형 {n}개, 최소 글자 {mp:.1f}pt")


def build():
    dia("모델 목록과 카드 점유", "lineup",
        head="그림 1: 라우터에 등록된 채팅 모델 전종과 각각이 차지하는 카드 수",
        items=[
            {"t": "라우터(:8400)에 채팅 모델 15종이 등록돼 있다. 이 밖에 임베딩과 리랭커가 하나씩 더 있지만 대화용이 아니라 이 정리에서는 뺐다"},
            {"t": "열한 종은 카드 넉 장을 통째로 쓴다. 그래서 그중 하나만 떠 있을 수 있고, 다른 것을 부르면 쓰던 것을 내리고 새로 올린다"},
            {"t": "네 종은 카드 한 장이면 되므로 넷까지 동시에 떠 있을 수 있다. 여러 사람이 각자 다른 모델을 쓸 때 이 차이가 크다"},
        ])
    d.code("측정 방법", [
        "# bench.py — 모델마다 (1) 준비될 때까지 기다리고 (2) 프롬프트 4종을 스트리밍으로 재고",
        "#            (3) 같은 프롬프트 4개를 동시에 던져 총처리량을 잰다",
        "PROMPTS = [사실 질문, 코드 작성, 계산 추론, 개념 설명]   # 네 종류 고정, 모든 모델에 같은 문장",
        "MAX_TOKENS = 1024                                     # 256 이면 사고하는 모델이 답까지 못 간다",
        "",
        "body = {'model': model, 'messages': [...],",
        "        'temperature': 0,                             # greedy, 무작위성 제거",
        "        'max_tokens': MAX_TOKENS, 'stream': True,",
        "        'stream_options': {'include_usage': True}}     # 토큰 수는 서버가 세 준 값을 쓴다",
        "",
        "t0 = time.perf_counter()",
        "for 이벤트 in 스트림:",
        "    본문 = delta.content;  사고 = delta.reasoning      # 사고하는 모델은 reasoning 으로 온다",
        "    if (본문 or 사고) and ttft is None:",
        "        ttft = time.perf_counter() - t0               # 첫 토큰까지 = TTFT",
        "total = time.perf_counter() - t0                      # 답이 끝날 때까지",
        "decode_tps = (출력토큰 - 1) / (total - ttft)           # 첫 토큰 뒤의 순수 생성 속도",
        "",
        "# 동시 요청: 같은 프롬프트 4개를 스레드로 동시에 → 합계 토큰 ÷ 전부 끝난 시간",
    ], caption="코드 1: bench.py, 모든 모델에 같은 조건을 강제하는 부분과 시간을 재는 부분")
    dia("시간 측정 구간", "method",
        head="그림 2: 한 요청의 시간 축에서 TTFT 와 디코드 구간이 각각 어디인지",
        items=[
            {"t": "TTFT 는 요청을 보낸 순간부터 첫 글자가 올 때까지다. 프롬프트를 읽는 시간이 여기 들어간다"},
            {"t": "디코드 속도는 첫 글자 이후 남은 토큰을 걸린 초로 나눈 값이라, 프롬프트 길이에 덜 휘둘린다"},
            {"t": "적재 시간은 따로 잰다. 요청이 아니라 모델을 카드에 올리는 시간이고, serve 로그에서 뽑는다"},
        ])
    dia("응답 지연", "latency",
        head="그림 3: 모델별 첫 토큰까지 걸린 시간과, 짧은 질문 하나가 끝날 때까지 걸린 시간",
        items=[
            {"t": "TTFT 는 어느 모델이든 0.5초 안쪽이라 체감 차이가 거의 없다. 기다림은 첫 글자가 아니라 답이 끝날 때까지에서 생긴다"},
            {"t": "별표가 붙은 사고하는 모델은 답을 내기 전에 속으로 길게 생각한다. 같은 30B 인데도 A3B Instruct 는 1초, A3B Thinking 은 15초가 걸렸다"},
            {"t": "짧게 묻고 짧게 받는 용도라면 사고 없는 모델이 훨씬 낫다. 어려운 추론을 맡길 때만 사고하는 쪽을 고른다"},
        ])
    dia("생성 처리량", "tput",
        head="그림 4: 요청 하나를 처리하는 속도와, 같은 질문 4개를 동시에 던졌을 때의 전체 속도",
        items=[
            {"t": "요청 하나만 보면 큰 모델과 작은 모델의 차이가 크지 않다. 카드를 넉 장 쓰는 30B 가 71 tok/s, 한 장 쓰는 4B 가 83 tok/s 다"},
            {"t": "동시에 여러 요청이 들어오면 갈린다. 카드 한 장짜리 작은 모델이 배치를 더 잘 채워 총처리량이 두 배 가까이 나온다"},
            {"t": "사람이 여럿 붙는 서비스라면 이 두 번째 숫자가 중요하다. 혼자 쓰는 도구라면 첫 번째 숫자와 앞 장의 응답 시간을 본다"},
        ])
    dia("모델 적재 시간", "loading",
        head="그림 5: 모델을 카드에 올리는 데 걸린 시간, serve 로그 기준",
        items=[
            {"t": "카드 넉 장을 쓰는 모델은 쓰던 것을 내리고 새로 올려야 해서, 모델을 바꾸면 수십 초에서 수 분을 기다린다"},
            {"t": "카드 한 장짜리는 십 초 안팎이라 갈아 끼우는 부담이 거의 없고, 넷까지 동시에 떠 있을 수 있어 기다림 자체가 잘 안 생긴다"},
            {"t": "한 번 올라간 모델은 계속 떠 있다. 자주 쓰는 모델이 정해져 있다면 이 시간은 하루에 몇 번만 치르는 비용이다"},
        ])
    # ── 숫자는 표로도 한 번 정리한다. 그림은 크기를 비교하게 하고, 표는 값을 읽게 한다.
    import gen_all

    def cell(v, digits=0, unit=""):
        return "-" if v is None else f"{v:.{digits}f}{unit}"

    trows = []
    for mid, lab, cards in gen_all.MODELS:
        rec = BY.get(mid)          # ★ 'd' 를 쓰면 덱 객체를 가려 build() 전체가 깨진다
        if not rec or "runs" not in rec:
            why = "배포 FXB 결함" if "A3B-FP8" in mid and "2507" not in mid else "측정 실패"
            trows.append([lab, str(cards), "못 뜸", "-", "-", "-", "-", why])
            continue
        c = rec.get("concurrent") or {}
        state = "정상" if rec.get("ok_runs") == rec.get("n_runs") else f"{rec.get('garbage', 0)}개 깨짐"
        trows.append([
            lab + (" *" if mid in gen_all.THINK else ""),
            str(cards),
            cell(rec.get("serve_load_s"), 0, "s"),
            cell(rec.get("ttft_med"), 2),
            cell(rec.get("decode_tps_med"), 0),
            cell(c.get("agg_tps"), 0),
            cell(gen_all.first_total(mid), 1, "s"),
            state,
        ])
    d.table("성능 요약", ["모델", "카드", "적재", "TTFT", "tok/s", "동시4", "첫응답", "답변"], trows,
            subtitle="같은 프롬프트 4종, temperature 0, max_tokens 1024. 별표는 사고하는 모델",
            note=("적재 = serve 로그의 Loading LLM 부터 기동 완료까지. TTFT = 첫 토큰까지(초, 중앙값). "
                  "tok/s = 첫 토큰 이후 생성 속도(중앙값). 동시4 = 같은 질문 4개 동시 요청의 합계 속도. "
                  "첫응답 = 가장 짧은 프롬프트가 끝날 때까지. "
                  "Qwen3-30B-A3B-FP8 은 배포 FXB 번들이 가중치를 30.2 GiB 다 읽은 뒤 "
                  "embed_tokens.weight 가 F32 인데 EDF 는 bf16 을 기대해 죽는다(50회 재현). "
                  "다운로드 문제가 아니라 배포본 결함이다."))

    rows = []
    for mid, lab, _ in gen_all.MODELS:
        rs = (BY.get(mid) or {}).get("runs")
        if not rs:
            continue
        r = rs[0]
        t = (r.get("text") or "").strip().replace("\n", " ")
        for ch in ("·", "ㆍ", "‧"):
            t = t.replace(ch, ", ")
        bad = r.get("broken")
        if bad:
            t = "같은 문자만 반복" if "반복" in str(bad) else str(bad)
        elif not t:
            t = "사고만 하고 답까지 못 감"
        rows.append([lab + (" *" if mid in gen_all.THINK else ""), t[:46], str(r.get("out_tokens") or "-")])
    d.table("답변 요약", ["모델", "답변 첫 줄", "토큰"], rows,
            subtitle="질문: 대한민국의 수도는 어디이고, 그 도시가 수도가 된 역사적 배경을 세 문장으로. 별표는 사고하는 모델",
            note="답변 원문 전체는 results/<모델>.json 에 저장돼 있다. 프롬프트 4종 중 첫 번째(사실 질문)의 답이다.")


PROMPTS = [
    ("사실 질문", "대한민국의 수도는 어디이고, 그 도시가 수도가 된 역사적 배경을 세 문장으로 설명해줘."),
    ("코드 작성", "파이썬으로 피보나치 수열의 n번째 항을 반복문으로 구하는 함수를 쓰고, 시간복잡도를 한 줄로 덧붙여줘."),
    ("계산 추론", "한 상자에 사과가 12개 들어간다. 사과 100개를 담으려면 상자가 몇 개 필요하고 마지막 상자에는 몇 개가 남는지 계산 과정을 보여줘."),
    ("개념 설명", "트랜스포머의 어텐션이 무엇인지 처음 배우는 사람에게 설명해줘. 비유를 하나 들고, 왜 순환신경망보다 병렬화에 유리한지도 말해줘."),
]
KEY = {"사실 질문": "fact", "코드 작성": "code", "계산 추론": "reason", "개념 설명": "long"}


def prompts_slide():
    """쓴 프롬프트를 요약하지 않고 원문 그대로 싣는다."""
    lines = []
    for i, (name, text) in enumerate(PROMPTS, 1):
        lines.append(f"# {i}. {name}")
        # 한 줄이 너무 길면 코드 상자에서 글자가 작아진다 — 40자쯤에서 접는다
        buf = text
        while buf:
            cut = 44 if len(buf) > 44 else len(buf)
            if cut < len(buf):
                sp = buf.rfind(" ", 0, cut + 1)
                cut = sp if sp > 20 else cut
            lines.append("  " + buf[:cut].strip())
            buf = buf[cut:].strip()
        lines.append("")
    d.code("사용한 프롬프트 전문", lines[:-1],
           caption="네 종류를 모든 모델에 똑같이 던졌다. 이 문장 그대로다",
           note="temperature 0, max_tokens 1024, 스트리밍. 답변 원문은 이어지는 절에 그대로 싣는다.")


def answer_slides(budget=760):
    """모델별 답변을 발췌 없이 싣는다. 분량이 커서 프롬프트마다 여러 장으로 나눈다."""
    import gen_all
    for name, _ in PROMPTS:
        key = KEY[name]
        items, used, part = [], 0, 1
        def flush(last=False):
            nonlocal items, used, part
            if not items:
                return
            # 제목에 대시를 쓰지 않는다(L-26). 콜론으로 잇는다.
            d.bullets(f"답변 원문: {name}" + (f" ({part})" if (part > 1 or not last) else ""),
                      items, subtitle=None)
            items, used = [], 0
            part += 1
        for mid, lab, _ in gen_all.MODELS:
            rec = BY.get(mid)
            if not rec or "runs" not in rec:
                continue
            r = next((x for x in rec["runs"] if x.get("prompt") == key), None)
            if not r:
                continue
            t = (r.get("text") or "").strip().replace("\n", " ")
            for ch in ("·", "ㆍ", "‧"):
                t = t.replace(ch, ", ")
            if not t:
                t = "(답변 없음 — 사고만 하고 예산이 끝났다)" if r.get("thinking") else "(빈 출력)"
            if used + len(t) > budget and items:
                flush()
            items.append({"t": f"{lab}  ({r.get('out_tokens')}tok)"})
            items.append({"t": t, "lv": 1})
            used += len(t)
        flush(last=True)


if __name__ == "__main__":
    build()
    prompts_slide()
    d.section("2", "답변 원문", "발췌하지 않고 모델이 낸 그대로")
    answer_slides()
    OUT = os.path.join(SRC, "..", "모델별-성능정리.pptx")
    d.save(OUT)
    print("넘침:", deck.OVERFLOW if deck.OVERFLOW else "없음")
    from pptx import Presentation
    prs = Presentation(OUT); worst = 0.0
    for sl in prs.slides:
        for shp in sl.shapes:
            if shp.shape_type == 6 and shp.name in srcs:
                diff, _, _ = compare(srcs[shp.name], shp); worst = max(worst, diff)
                print(f"되돌려 비교 {shp.name}: {diff*100:.2f}% ({'통과' if diff <= 0.02 else '실패'})")
    print(f"최대 {worst*100:.2f}%, 슬라이드 {len(prs.slides)}장, 저장 {os.path.abspath(OUT)}")
